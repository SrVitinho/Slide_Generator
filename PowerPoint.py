import pptx
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from icrawler.builtin import GoogleImageCrawler
from pathlib import Path
import os
import PPRT_Aux as aux


class Ppt:
    def __init__(self, title, resumos):
        self.title = title
        self.resumos = resumos
        self.prs = pptx.Presentation()
        self.image_slide_layout = self.prs.slide_layouts[5]
        self.page = 0
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 32, 96)
        title.text = self.title
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    def add_page(self, title):
        slide = self.prs.slides.add_slide(self.image_slide_layout)
        shapes = slide.shapes

        # Image - Background
        left = top = Inches(0)
        filters = dict(
            size='=1980x1080')

        current_path = '/home/ec2-user/script/'
        img_folder = current_path + 'Images/'
        background_path = img_folder

        list = os.listdir(img_folder)
        try:
            for picture in list:
                remove_path = img_folder + picture
                os.remove(remove_path)
        except Exception as err:
            print(err)
        else:
            print('Image removed')

        google_crawler = GoogleImageCrawler(storage={'root_dir': background_path})
        google_crawler.crawl(keyword=title, max_num=1, filters=filters)
        img_folder_crawler = img_folder
        list = os.listdir(img_folder_crawler)
        img_path = img_folder_crawler + list[0]

        pic = slide.shapes.add_picture(img_path, left, top, width=self.prs.slide_width, height=self.prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Create Title
        title_shape = shapes.title
        title_shape.text = title
        '''# Add Image
        left = Inches(0.75)
        top = Inches(1.25)
        pic = slide.shapes.add_picture(image, left, top, height=Inches(4.5))'''

        # Build the Text Box
        print('ppt resumo pre:' + self.resumos[self.page])
        print('------------------------------------------')
        clean_resumo = aux.fix_lines(self.resumos[self.page])
        print('Clean: ' + clean_resumo)
        print('---------------------------------------------')
        top, height = aux.get_sizes(clean_resumo)
        left = Inches(0.75)
        width = Inches(9)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        #Prepare Text
        list_resumo = aux.get_text_pages(clean_resumo)

        # Title
        p = tf.add_paragraph()
        print(list_resumo)
        p.text = list_resumo[0]
        p.font.size = Pt(18)

        # Add Pages
        txBox = slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.font.size = Pt(15)
        self.page += 1

        for clean_list_resumo in list_resumo[1:]:
            self.sub_page(title, clean_list_resumo)

    def sub_page(self, title, text):
        slide = self.prs.slides.add_slide(self.image_slide_layout)
        shapes = slide.shapes

        # Image - Background
        left = top = Inches(0)

        current_path = '/home/ec2-user/script/'
        img_folder = current_path + 'Images/'

        img_folder_crawler = img_folder
        list = os.listdir(img_folder_crawler)
        img_path = img_folder_crawler + list[0]

        pic = slide.shapes.add_picture(img_path, left, top, width=self.prs.slide_width, height=self.prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Create Title
        title_shape = shapes.title
        title_shape.text = title

        # Build the Text Box
        top, height = aux.get_sizes(text)
        left = Inches(0.75)
        width = Inches(9)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        # Title
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(18)

        # Add Pages
        txBox = slide.shapes.add_textbox(Inches(9), Inches(6.75), Inches(1), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.font.size = Pt(15)

    def save(self):
        self.prs.save('/home/ec2-user/output/presentation.pptx')
